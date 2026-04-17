#!/usr/bin/env python3
"""
Character Sheet Presentation Generator
========================================
Creates a 16:9 PowerPoint presentation from rendered character images.

Usage:
  python make_presentation.py --image ref.png --name "Character Name" --desc "Description" --output-dir ./renders
  python make_presentation.py --image ref.png --name "Nora" --desc "A curious 10-year-old adventurer" --output-dir ./pose_full_output
"""

import argparse
import os
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import io


# Slide dimensions: 16:9
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
ACCENT_COLOR = RGBColor(0x33, 0x66, 0x99)
SUBTLE_COLOR = RGBColor(0x77, 0x77, 0x77)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_image_fitted(slide, img_path, left, top, max_width, max_height):
    """Add image maintaining aspect ratio within bounds."""
    img = Image.open(img_path)
    w, h = img.size
    scale = min(max_width / w, max_height / h)
    final_w = int(w * scale)
    final_h = int(h * scale)
    # Center within bounds
    offset_x = (max_width - final_w) // 2
    offset_y = (max_height - final_h) // 2
    slide.shapes.add_picture(img_path, left + offset_x, top + offset_y,
                             Emu(final_w), Emu(final_h))


def find_images(output_dir):
    """Find rendered images, skeletons, and organize by type."""
    renders = sorted(glob.glob(os.path.join(output_dir, "*.png")))
    renders = [r for r in renders if "grid" not in os.path.basename(r)]

    skeletons = sorted(glob.glob(os.path.join(output_dir, "poses", "*_skeleton.png")))

    # Categorize renders
    front_views = [r for r in renders if "front_view" in r and "eyelevel" in r and "d1.0" in r]
    side_views = [r for r in renders if "side_view" in r and "eyelevel" in r]
    quarter_views = [r for r in renders if "quarter_view" in r and "eyelevel" in r]
    elevated = [r for r in renders if "elevated" in r or "highangle" in r]
    low_angle = [r for r in renders if "lowangle" in r]

    # Try selects first
    selects_dir = os.path.join(output_dir, "selects")
    if os.path.isdir(selects_dir):
        selects = sorted(glob.glob(os.path.join(selects_dir, "*.png")))
        if selects:
            return {
                "all": renders,
                "selects": selects,
                "front": front_views,
                "side": side_views,
                "quarter": quarter_views,
                "elevated": elevated,
                "low_angle": low_angle,
                "skeletons": skeletons,
            }

    return {
        "all": renders,
        "selects": [],
        "front": front_views,
        "side": side_views,
        "quarter": quarter_views,
        "elevated": elevated,
        "low_angle": low_angle,
        "skeletons": skeletons,
    }


def make_title_slide(prs, ref_image, char_name, char_desc):
    """Page 1: A-pose reference on left, name and description on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_COLOR)

    # Left column: reference image
    img_left = Inches(0.8)
    img_top = Inches(0.6)
    img_max_w = Emu(int(SLIDE_WIDTH / 2 - Inches(1.2)))
    img_max_h = Emu(int(SLIDE_HEIGHT - Inches(1.2)))
    add_image_fitted(slide, ref_image, img_left, img_top, img_max_w, img_max_h)

    # Right column: name + description
    text_left = Emu(int(SLIDE_WIDTH / 2 + Inches(0.4)))
    text_width = Emu(int(SLIDE_WIDTH / 2 - Inches(1.2)))

    # Character name
    add_text(slide, text_left, Inches(2.0), text_width, Inches(1.2),
             char_name, font_size=48, color=TEXT_COLOR, bold=True)

    # Divider line
    add_text(slide, text_left, Inches(3.2), text_width, Inches(0.1),
             "─" * 30, font_size=12, color=ACCENT_COLOR)

    # Description
    add_text(slide, text_left, Inches(3.6), text_width, Inches(2.5),
             char_desc, font_size=20, color=SUBTLE_COLOR)

    # Footer
    add_text(slide, text_left, Inches(6.2), text_width, Inches(0.5),
             "Character Sheet", font_size=14, color=ACCENT_COLOR)


def make_angles_slide(prs, images, title="Multi-Angle Views"):
    """Page 2: Grid of angle variations — 4x2 layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             title, font_size=28, color=TEXT_COLOR, bold=True)

    # Subtitle line
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
             "Camera orbit variations from a single reference image",
             font_size=14, color=SUBTLE_COLOR)

    # Grid: 4 columns x 2 rows
    picks = images[:8]
    cols, rows = 4, 2
    margin_x = Inches(0.6)
    margin_top = Inches(1.4)
    spacing = Inches(0.15)
    avail_w = int(SLIDE_WIDTH - margin_x * 2 - spacing * (cols - 1))
    avail_h = int(SLIDE_HEIGHT - margin_top - Inches(0.5) - spacing * (rows - 1))
    cell_w = avail_w // cols
    cell_h = avail_h // rows

    for i, img_path in enumerate(picks):
        r, c = i // cols, i % cols
        left = int(margin_x + c * (cell_w + spacing))
        top = int(margin_top + r * (cell_h + spacing))
        add_image_fitted(slide, img_path, left, top, cell_w, cell_h)

        # Small label under each image
        basename = os.path.splitext(os.path.basename(img_path))[0]
        # Extract human-readable angle from filename
        parts = basename.split("_")
        label = " ".join(p for p in parts if not p.startswith("az") and not p.startswith("el") and not p.startswith("d"))
        label = label.replace("shot ", "").strip()
        if len(label) > 30:
            label = label[:30]
        add_text(slide, left, top + cell_h - Inches(0.05), cell_w, Inches(0.35),
                 label, font_size=9, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)


def make_detail_slide(prs, images, skeletons, title="Poses & Skeleton Analysis"):
    """Page 3: Mix of renders and skeletons in a 3x2 layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    # Title
    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             title, font_size=28, color=TEXT_COLOR, bold=True)

    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
             "Rendered views paired with DWPose skeleton extraction",
             font_size=14, color=SUBTLE_COLOR)

    # 3 columns x 2 rows: top row = renders, bottom row = matching skeletons
    picks = images[:3]
    cols = 3
    margin_x = Inches(0.8)
    margin_top = Inches(1.4)
    spacing_x = Inches(0.3)
    spacing_y = Inches(0.2)
    avail_w = int(SLIDE_WIDTH - margin_x * 2 - spacing_x * (cols - 1))
    cell_w = avail_w // cols
    row_h = int((SLIDE_HEIGHT - margin_top - Inches(0.4) - spacing_y) / 2)

    for i, img_path in enumerate(picks):
        c = i
        left = int(margin_x + c * (cell_w + spacing_x))

        # Top row: render
        add_image_fitted(slide, img_path, left, int(margin_top), cell_w, row_h)

        # Bottom row: matching skeleton
        basename = os.path.splitext(os.path.basename(img_path))[0]
        skel_path = None
        for s in skeletons:
            if basename in os.path.basename(s):
                skel_path = s
                break

        if skel_path and os.path.exists(skel_path):
            add_image_fitted(slide, skel_path, left,
                           int(margin_top + row_h + spacing_y), cell_w, row_h)

        # Label
        parts = basename.split("_")
        label = " ".join(p for p in parts if not p.startswith("az") and not p.startswith("el") and not p.startswith("d"))
        label = label.replace("shot ", "").strip()
        add_text(slide, left, int(margin_top) - Inches(0.25), cell_w, Inches(0.25),
                 label, font_size=10, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)


def make_expressions_slide(prs, images):
    """Page 4: 4x4 grid of facial expressions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Expressions", font_size=28, color=TEXT_COLOR, bold=True)

    add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
             "Facial expression variations driven by prompt editing",
             font_size=14, color=SUBTLE_COLOR)

    picks = images[:16]
    cols, rows = 4, 4
    margin_x = Inches(0.6)
    margin_top = Inches(1.3)
    spacing = Inches(0.12)
    avail_w = int(SLIDE_WIDTH - margin_x * 2 - spacing * (cols - 1))
    avail_h = int(SLIDE_HEIGHT - margin_top - Inches(0.3) - spacing * (rows - 1))
    cell_w = avail_w // cols
    cell_h = avail_h // rows

    for i, img_path in enumerate(picks):
        r, c = i // cols, i % cols
        left = int(margin_x + c * (cell_w + spacing))
        top = int(margin_top + r * (cell_h + spacing))
        add_image_fitted(slide, img_path, left, top, cell_w, cell_h)

        # Label from filename
        basename = os.path.splitext(os.path.basename(img_path))[0]
        label = basename.replace("expr_", "").replace("_", " ").title()
        add_text(slide, left, top + cell_h - Inches(0.02), cell_w, Inches(0.3),
                 label, font_size=9, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)


def make_outfits_lighting_slide(prs, outfits, lighting):
    """Page 5: Outfits on top row, lighting on bottom row."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)

    add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
             "Outfits & Lighting", font_size=28, color=TEXT_COLOR, bold=True)

    margin_x = Inches(0.8)
    spacing_x = Inches(0.25)

    # Top row: Outfits
    add_text(slide, Inches(0.6), Inches(1.0), Inches(3), Inches(0.35),
             "Outfits", font_size=16, color=ACCENT_COLOR)

    outfit_picks = outfits[:4]
    cols = len(outfit_picks) if outfit_picks else 4
    row_top = Inches(1.4)
    row_h = int((SLIDE_HEIGHT - Inches(3.2)) / 2)
    avail_w = int(SLIDE_WIDTH - margin_x * 2 - spacing_x * (cols - 1))
    cell_w = avail_w // cols if cols > 0 else avail_w

    for i, img_path in enumerate(outfit_picks):
        left = int(margin_x + i * (cell_w + spacing_x))
        add_image_fitted(slide, img_path, left, int(row_top), cell_w, row_h)
        basename = os.path.splitext(os.path.basename(img_path))[0]
        label = basename.replace("outfit_", "").replace("_", " ").title()
        add_text(slide, left, int(row_top) + row_h, cell_w, Inches(0.3),
                 label, font_size=10, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)

    # Bottom row: Lighting
    row2_top = int(row_top) + row_h + Inches(0.6)
    add_text(slide, Inches(0.6), row2_top - Inches(0.35), Inches(3), Inches(0.35),
             "Lighting", font_size=16, color=ACCENT_COLOR)

    light_picks = lighting[:4]
    cols = len(light_picks) if light_picks else 4

    for i, img_path in enumerate(light_picks):
        left = int(margin_x + i * (cell_w + spacing_x))
        add_image_fitted(slide, img_path, left, int(row2_top), cell_w, row_h)
        basename = os.path.splitext(os.path.basename(img_path))[0]
        label = basename.replace("light_", "").replace("_", " ").title()
        add_text(slide, left, int(row2_top) + row_h, cell_w, Inches(0.3),
                 label, font_size=10, color=SUBTLE_COLOR, alignment=PP_ALIGN.CENTER)


def generate_presentation(ref_image, char_name, char_desc, output_dir, output_file=None,
                          expressions_dir=None, outfits_dir=None, lighting_dir=None):
    """Generate the full character sheet presentation."""
    if output_file is None:
        output_file = os.path.join(output_dir, f"{char_name.lower().replace(' ', '_')}_character_sheet.pptx")

    images = find_images(output_dir)

    # Use selects if available, otherwise pick from all renders
    angle_images = images["selects"] if images["selects"] else images["all"][:16]
    skeleton_images = images["skeletons"]

    # Pick best images for each slide
    # Slide 2: mix of angles
    slide2_picks = angle_images[:8]

    # Slide 3: 3 renders with matching skeletons — pick ones that have skeletons
    slide3_picks = []
    for img in angle_images:
        basename = os.path.splitext(os.path.basename(img))[0]
        has_skel = any(basename in os.path.basename(s) for s in skeleton_images)
        if has_skel:
            slide3_picks.append(img)
        if len(slide3_picks) >= 3:
            break
    if len(slide3_picks) < 3:
        slide3_picks = angle_images[:3]

    # Build presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Page 1: Title with A-pose
    make_title_slide(prs, ref_image, char_name, char_desc)

    # Page 2: Angle grid
    make_angles_slide(prs, slide2_picks)

    # Page 3: Renders + skeletons
    make_detail_slide(prs, slide3_picks, skeleton_images)

    # Page 4: Expressions
    base_dir = os.path.dirname(os.path.abspath(output_dir))
    expr_search = [expressions_dir] if expressions_dir else glob.glob(os.path.join(base_dir, "*expressions*"))
    expr_images = []
    for d in expr_search:
        if d and os.path.isdir(d):
            expr_images.extend(sorted(glob.glob(os.path.join(d, "expr_*.png"))))
    if expr_images:
        make_expressions_slide(prs, expr_images)

    # Page 5: Outfits & Lighting
    outfit_search = [outfits_dir] if outfits_dir else glob.glob(os.path.join(base_dir, "*outfits*"))
    outfit_images = []
    for d in outfit_search:
        if d and os.path.isdir(d):
            outfit_images.extend(sorted(glob.glob(os.path.join(d, "outfit_*.png"))))

    light_search = [lighting_dir] if lighting_dir else glob.glob(os.path.join(base_dir, "*lighting*"))
    light_images = []
    for d in light_search:
        if d and os.path.isdir(d):
            light_images.extend(sorted(glob.glob(os.path.join(d, "light_*.png"))))

    if outfit_images or light_images:
        make_outfits_lighting_slide(prs, outfit_images, light_images)

    prs.save(output_file)
    print(f"Saved presentation: {output_file}")
    return output_file


def main():
    parser = argparse.ArgumentParser(description="Generate character sheet PowerPoint presentation")
    parser.add_argument("--image", required=True, help="Reference A-pose image")
    parser.add_argument("--name", required=True, help="Character name")
    parser.add_argument("--desc", default="", help="Character description")
    parser.add_argument("--output-dir", required=True, help="Directory with rendered images")
    parser.add_argument("--output", default=None, help="Output .pptx path")
    parser.add_argument("--expressions-dir", default=None, help="Directory with expression images")
    parser.add_argument("--outfits-dir", default=None, help="Directory with outfit images")
    parser.add_argument("--lighting-dir", default=None, help="Directory with lighting images")
    args = parser.parse_args()

    generate_presentation(args.image, args.name, args.desc, args.output_dir, args.output,
                          args.expressions_dir, args.outfits_dir, args.lighting_dir)


if __name__ == "__main__":
    main()
